import json
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from computer_agent.documents import document_support_status, read_structured_document
from computer_agent.runtime_tools import TOOL_SCHEMAS, LocalOnlyResult, execute_tool


def _write_test_pdf(path: Path) -> None:
    stream = b"\n".join(
        [
            b"BT /F1 18 Tf 72 720 Td (Quarterly Report) Tj ET",
            b"BT /F1 11 Tf 72 690 Td (Revenue increased this quarter.) Tj ET",
            b"0.5 w",
            b"72 640 m 300 640 l S",
            b"72 620 m 300 620 l S",
            b"72 600 m 300 600 l S",
            b"72 600 m 72 640 l S",
            b"180 600 m 180 640 l S",
            b"300 600 m 300 640 l S",
            b"BT /F1 10 Tf 80 627 Td (Item) Tj ET",
            b"BT /F1 10 Tf 188 627 Td (Value) Tj ET",
            b"BT /F1 10 Tf 80 607 Td (Revenue) Tj ET",
            b"BT /F1 10 Tf 188 607 Td (120) Tj ET",
        ]
    )
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
    ]

    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{index} 0 obj\n".encode())
        pdf.extend(obj)
        pdf.extend(b"\nendobj\n")
    xref = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode())
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n"
        ).encode()
    )
    path.write_bytes(pdf)


def test_document_support_status_lists_all_backends() -> None:
    assert document_support_status() == "Document support: PDF DOCX XLSX XLSM PPTX"


def test_pdf_reader_extracts_pages_headings_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "report.pdf"
    _write_test_pdf(path)

    result = json.loads(read_structured_document(str(path)))

    assert result["document_type"] == "pdf"
    assert result["page_count"] == 1
    page = result["pages"][0]
    assert any("Quarterly Report" in text for text in page["paragraphs"])
    assert any(item["text"] == "Quarterly Report" for item in page["headings"])
    assert "tables" in page
    assert result["read_only"] is True


def test_word_reader_preserves_heading_paragraph_and_table_order(tmp_path: Path) -> None:
    path = tmp_path / "brief.docx"
    document = Document()
    document.core_properties.title = "Client Brief"
    document.add_heading("Scope", level=1)
    document.add_paragraph("Build the document layer.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Interly"
    table.cell(1, 1).text = "Active"
    document.save(path)

    result = json.loads(read_structured_document(str(path)))

    assert result["document_type"] == "word"
    assert result["metadata"]["title"] == "Client Brief"
    assert result["blocks"][0] == {"kind": "heading", "level": 1, "text": "Scope"}
    assert result["blocks"][1]["text"] == "Build the document layer."
    table_block = next(block for block in result["blocks"] if block["kind"] == "table")
    assert table_block["rows"][1] == ["Interly", "Active"]


def test_excel_reader_exposes_sheet_headers_range_and_formulas(tmp_path: Path) -> None:
    path = tmp_path / "sales.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Item", "Qty", "Total"])
    sheet.append(["Widget", 2, "=B2*10"])
    workbook.save(path)

    result = json.loads(read_structured_document(str(path)))

    assert result["document_type"] == "excel"
    assert result["sheet_count"] == 1
    parsed_sheet = result["sheets"][0]
    assert parsed_sheet["name"] == "Sales"
    assert parsed_sheet["range"] == "A1:C2"
    assert parsed_sheet["headers"] == ["Item", "Qty", "Total"]
    assert parsed_sheet["formula_cells"] == [{"cell": "C2", "formula": "=B2*10"}]


def test_powerpoint_reader_exposes_slide_text_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "roadmap.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Priority 3"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    textbox.text_frame.text = "Documents become understandable."
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(1.5))
    table_shape.table.cell(0, 0).text = "Format"
    table_shape.table.cell(0, 1).text = "State"
    table_shape.table.cell(1, 0).text = "PDF"
    table_shape.table.cell(1, 1).text = "Ready"
    presentation.save(path)

    result = json.loads(read_structured_document(str(path)))

    assert result["document_type"] == "powerpoint"
    assert result["slide_count"] == 1
    parsed_slide = result["slides"][0]
    assert parsed_slide["title"] == "Priority 3"
    assert parsed_slide["blocks"][0]["paragraphs"][0]["text"] == "Documents become understandable."
    assert parsed_slide["tables"][0]["rows"][1] == ["PDF", "Ready"]


def test_structured_read_reuses_sensitive_read_text_file_tool(tmp_path: Path) -> None:
    path = tmp_path / "memory.xlsx"
    workbook = Workbook()
    workbook.active["A1"] = "private value"
    workbook.save(path)

    result = execute_tool("read_text_file", json.dumps({"path": str(path)}))

    assert isinstance(result, LocalOnlyResult)
    assert "private value" in result.terminal_output
    assert "private value" not in result.model_status
    schema = next(
        item["function"]
        for item in TOOL_SCHEMAS
        if item["function"]["name"] == "read_text_file"
    )
    assert "PDF" in schema["description"]
    assert "PowerPoint" in schema["description"]


def test_corrupt_office_package_fails_without_parsing(tmp_path: Path) -> None:
    path = tmp_path / "bad.docx"
    path.write_text("not an Office package", encoding="utf-8")

    result = read_structured_document(str(path))

    assert "encrypted, corrupt, or not a valid Open XML package" in result
