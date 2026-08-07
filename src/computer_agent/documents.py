"""Read PDF and Office documents into one bounded, read-only structure."""

from __future__ import annotations

import json
import re
import zipfile
from dataclasses import dataclass
from datetime import date, datetime, time
from pathlib import Path
from statistics import median
from typing import Any

import pdfplumber
from docx import Document as WordDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation

SUPPORTED_EXTENSIONS = {
    ".pdf": "pdf",
    ".docx": "word",
    ".xlsx": "excel",
    ".xlsm": "excel",
    ".pptx": "powerpoint",
}
MAX_DOCUMENT_BYTES = 50 * 1024 * 1024
MAX_OFFICE_UNCOMPRESSED_BYTES = 250 * 1024 * 1024
MAX_OFFICE_ZIP_ENTRIES = 10_000
MAX_TOTAL_TEXT_CHARS = 90_000
MAX_TEXT_ITEM_CHARS = 8_000
MAX_PAGES = 200
MAX_SLIDES = 300
MAX_SHEETS = 50
MAX_TABLE_ROWS = 200
MAX_TABLE_COLUMNS = 50
MAX_BLOCKS = 2_000


@dataclass
class OutputBudget:
    """Bound document output before it is exposed to the terminal or model."""

    remaining_chars: int = MAX_TOTAL_TEXT_CHARS
    truncated: bool = False

    def take(self, value: Any, limit: int = MAX_TEXT_ITEM_CHARS) -> str:
        text = "" if value is None else str(value)
        if not text:
            return ""
        allowed = min(limit, self.remaining_chars)
        if allowed <= 0:
            self.truncated = True
            return ""
        if len(text) > allowed:
            self.truncated = True
            suffix = "\n[truncated]"
            if allowed > len(suffix):
                text = text[: allowed - len(suffix)] + suffix
            else:
                text = text[:allowed]
        self.remaining_chars -= len(text)
        return text


def document_support_status() -> str:
    """Return a frozen-build smoke-test string after all document backends import."""
    backends = [
        pdfplumber.__name__,
        WordDocument.__module__,
        load_workbook.__module__,
        Presentation.__module__,
    ]
    if not all(backends):
        raise RuntimeError("One or more document backends failed to import.")
    return "Document support: PDF DOCX XLSX XLSM PPTX"


def read_structured_document(path: str) -> str:
    """Read one supported document into a common bounded JSON representation."""
    target = Path(path).expanduser()
    if not target.is_file():
        return f"Document not found: {target}"
    suffix = target.suffix.casefold()
    kind = SUPPORTED_EXTENSIONS.get(suffix)
    if kind is None:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        return f"Unsupported document type {suffix or '(no extension)'}. Supported: {supported}"
    try:
        size = target.stat().st_size
    except OSError as error:
        return f"Could not inspect document: {error}"
    if size > MAX_DOCUMENT_BYTES:
        return (
            f"Document is too large ({size} bytes). "
            f"Interly limits structured document reads to {MAX_DOCUMENT_BYTES} bytes."
        )
    if suffix in {".docx", ".xlsx", ".xlsm", ".pptx"}:
        validation_error = _validate_office_container(target)
        if validation_error:
            return validation_error

    budget = OutputBudget()
    try:
        if kind == "pdf":
            result = _read_pdf(target, budget)
        elif kind == "word":
            result = _read_word(target, budget)
        elif kind == "excel":
            result = _read_excel(target, budget)
        else:
            result = _read_powerpoint(target, budget)
    except (OSError, ValueError, KeyError, TypeError, zipfile.BadZipFile) as error:
        return f"Structured document read failed safely: {error}"

    result["path"] = str(target.resolve())
    result["size_bytes"] = size
    result["read_only"] = True
    result["truncated"] = bool(result.get("truncated") or budget.truncated)
    result["limits"] = {
        "max_document_bytes": MAX_DOCUMENT_BYTES,
        "max_total_text_chars": MAX_TOTAL_TEXT_CHARS,
        "max_table_rows": MAX_TABLE_ROWS,
        "max_table_columns": MAX_TABLE_COLUMNS,
    }
    return json.dumps(result, indent=2, ensure_ascii=False, default=_json_default)


def _validate_office_container(path: Path) -> str | None:
    try:
        with zipfile.ZipFile(path) as archive:
            members = archive.infolist()
    except zipfile.BadZipFile:
        return "Office document is encrypted, corrupt, or not a valid Open XML package."
    if len(members) > MAX_OFFICE_ZIP_ENTRIES:
        return "Office document contains too many package entries; nothing was read."
    uncompressed = sum(member.file_size for member in members)
    if uncompressed > MAX_OFFICE_UNCOMPRESSED_BYTES:
        return (
            "Office document expands beyond Interly's safe read limit "
            f"({uncompressed} bytes uncompressed); nothing was read."
        )
    return None


def _read_pdf(path: Path, budget: OutputBudget) -> dict[str, Any]:
    pages: list[dict[str, Any]] = []
    with pdfplumber.open(path) as document:
        metadata = {
            str(key): budget.take(value, 1_000)
            for key, value in (document.metadata or {}).items()
            if value not in {None, ""}
        }
        page_count = len(document.pages)
        for index, page in enumerate(document.pages[:MAX_PAGES], start=1):
            page_text = page.extract_text() or ""
            paragraphs = [
                budget.take(block.strip())
                for block in re.split(r"\n\s*\n", page_text)
                if block.strip()
            ]
            if len(paragraphs) <= 1 and page_text:
                paragraphs = [
                    budget.take(line.strip())
                    for line in page_text.splitlines()
                    if line.strip()
                ]
            headings = _pdf_headings(page, budget)
            tables = _bounded_tables(page.extract_tables() or [], budget)
            pages.append(
                {
                    "kind": "page",
                    "index": index,
                    "width": round(float(page.width), 2),
                    "height": round(float(page.height), 2),
                    "headings": headings,
                    "paragraphs": paragraphs[:MAX_BLOCKS],
                    "tables": tables,
                }
            )
            if budget.remaining_chars <= 0:
                break
    return {
        "document_type": "pdf",
        "metadata": metadata,
        "page_count": page_count,
        "pages": pages,
        "truncated": page_count > len(pages),
    }


def _pdf_headings(page: Any, budget: OutputBudget) -> list[dict[str, Any]]:
    words = page.extract_words(extra_attrs=["size"], use_text_flow=True) or []
    sizes = [float(word.get("size", 0) or 0) for word in words if word.get("size")]
    if not sizes:
        return []
    body_size = median(sizes)
    lines: list[list[dict[str, Any]]] = []
    sorted_words = sorted(
        words,
        key=lambda item: (float(item.get("top", 0)), float(item.get("x0", 0))),
    )
    for word in sorted_words:
        top = float(word.get("top", 0))
        if not lines or abs(top - float(lines[-1][0].get("top", 0))) > 3:
            lines.append([word])
        else:
            lines[-1].append(word)
    headings: list[dict[str, Any]] = []
    for line in lines:
        text = " ".join(str(word.get("text", "")) for word in line).strip()
        line_size = max(float(word.get("size", 0) or 0) for word in line)
        if text and len(text) <= 180 and line_size >= max(body_size + 1.0, body_size * 1.12):
            headings.append(
                {
                    "text": budget.take(text, 500),
                    "font_size": round(line_size, 2),
                    "top": round(float(line[0].get("top", 0)), 2),
                }
            )
        if len(headings) >= 100 or budget.remaining_chars <= 0:
            break
    return headings


def _read_word(path: Path, budget: OutputBudget) -> dict[str, Any]:
    document = WordDocument(path)
    blocks: list[dict[str, Any]] = []
    for child in document.element.body.iterchildren():
        if len(blocks) >= MAX_BLOCKS or budget.remaining_chars <= 0:
            budget.truncated = True
            break
        if isinstance(child, CT_P):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if not text:
                continue
            style = paragraph.style.name if paragraph.style is not None else ""
            heading_match = re.match(r"Heading\s+(\d+)", style or "", re.IGNORECASE)
            if heading_match:
                blocks.append(
                    {
                        "kind": "heading",
                        "level": int(heading_match.group(1)),
                        "text": budget.take(text),
                    }
                )
            else:
                blocks.append(
                    {"kind": "paragraph", "style": style or None, "text": budget.take(text)}
                )
        elif isinstance(child, CT_Tbl):
            table = Table(child, document)
            table_rows = list(table.rows)
            rows = []
            for row in table_rows[:MAX_TABLE_ROWS]:
                cells = list(row.cells)
                rows.append(
                    [
                        budget.take(cell.text.strip(), 2_000)
                        for cell in cells[:MAX_TABLE_COLUMNS]
                    ]
                )
            blocks.append(
                {
                    "kind": "table",
                    "rows": rows,
                    "row_count": len(table_rows),
                    "column_count": max(
                        (len(list(row.cells)) for row in table_rows),
                        default=0,
                    ),
                    "truncated": len(table_rows) > MAX_TABLE_ROWS
                    or any(len(list(row.cells)) > MAX_TABLE_COLUMNS for row in table_rows),
                }
            )
    core = document.core_properties
    metadata = {
        "title": budget.take(core.title, 1_000),
        "subject": budget.take(core.subject, 1_000),
        "author": budget.take(core.author, 1_000),
        "keywords": budget.take(core.keywords, 1_000),
        "comments": budget.take(core.comments, 2_000),
        "created": _json_default(core.created) if core.created else None,
        "modified": _json_default(core.modified) if core.modified else None,
    }
    return {
        "document_type": "word",
        "metadata": {
            key: value for key, value in metadata.items() if value not in {None, ""}
        },
        "blocks": blocks,
        "paragraph_count": len(document.paragraphs),
        "table_count": len(document.tables),
        "truncated": budget.truncated,
    }


def _read_excel(path: Path, budget: OutputBudget) -> dict[str, Any]:
    workbook = load_workbook(path, read_only=True, data_only=False, keep_links=False)
    sheets: list[dict[str, Any]] = []
    try:
        sheet_names = workbook.sheetnames
        for sheet_name in sheet_names[:MAX_SHEETS]:
            if budget.remaining_chars <= 0:
                budget.truncated = True
                break
            sheet = workbook[sheet_name]
            captured_rows: list[list[Any]] = []
            formula_cells: list[dict[str, Any]] = []
            header_row: list[Any] | None = None
            row_count = 0
            max_columns_seen = 0
            for row in sheet.iter_rows():
                row_count += 1
                values: list[Any] = []
                for cell in row[:MAX_TABLE_COLUMNS]:
                    value = _excel_value(cell.value, budget)
                    values.append(value)
                    if (
                        isinstance(cell.value, str)
                        and cell.value.startswith("=")
                        and len(formula_cells) < 500
                    ):
                        formula_cells.append(
                            {
                                "cell": cell.coordinate,
                                "formula": budget.take(cell.value, 2_000),
                            }
                        )
                max_columns_seen = max(max_columns_seen, len(row))
                if header_row is None and any(value not in {None, ""} for value in values):
                    header_row = values
                if len(captured_rows) < MAX_TABLE_ROWS:
                    captured_rows.append(values)
                else:
                    budget.truncated = True
                if row_count >= 5_000:
                    budget.truncated = True
                    break
                if budget.remaining_chars <= 0:
                    break
            max_row = int(getattr(sheet, "max_row", row_count) or row_count)
            max_column = int(getattr(sheet, "max_column", max_columns_seen) or max_columns_seen)
            dimension = (
                f"A1:{get_column_letter(max(1, max_column))}{max(1, max_row)}"
                if max_row or max_column
                else "A1"
            )
            sheets.append(
                {
                    "kind": "sheet",
                    "name": budget.take(sheet_name, 500),
                    "state": str(getattr(sheet, "sheet_state", "visible")),
                    "range": dimension,
                    "row_count": max_row,
                    "column_count": max_column,
                    "headers": header_row or [],
                    "rows": captured_rows,
                    "formula_cells": formula_cells,
                    "truncated": max_row > len(captured_rows)
                    or max_column > MAX_TABLE_COLUMNS,
                }
            )
    finally:
        workbook.close()
    return {
        "document_type": "excel",
        "sheet_count": len(sheet_names),
        "sheets": sheets,
        "truncated": len(sheet_names) > len(sheets) or budget.truncated,
    }


def _read_powerpoint(path: Path, budget: OutputBudget) -> dict[str, Any]:
    presentation = Presentation(path)
    slides: list[dict[str, Any]] = []
    slide_list = list(presentation.slides)
    for index, slide in enumerate(slide_list[:MAX_SLIDES], start=1):
        if budget.remaining_chars <= 0:
            budget.truncated = True
            break
        title_shape = slide.shapes.title
        title = (
            budget.take(title_shape.text.strip(), 1_000)
            if title_shape is not None and title_shape.text
            else ""
        )
        blocks: list[dict[str, Any]] = []
        tables: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table_rows = list(shape.table.rows)
                rows = []
                for row in table_rows[:MAX_TABLE_ROWS]:
                    cells = list(row.cells)
                    rows.append(
                        [
                            budget.take(cell.text.strip(), 2_000)
                            for cell in cells[:MAX_TABLE_COLUMNS]
                        ]
                    )
                tables.append(
                    {
                        "rows": rows,
                        "row_count": len(table_rows),
                        "column_count": len(shape.table.columns),
                        "truncated": len(table_rows) > MAX_TABLE_ROWS
                        or len(shape.table.columns) > MAX_TABLE_COLUMNS,
                    }
                )
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            if title_shape is not None and shape == title_shape:
                continue
            paragraphs = []
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    paragraphs.append(
                        {"level": int(paragraph.level), "text": budget.take(text, 4_000)}
                    )
            if paragraphs:
                blocks.append({"kind": "text", "paragraphs": paragraphs})
        notes = ""
        if slide.has_notes_slide:
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame is not None:
                notes = budget.take((notes_frame.text or "").strip(), 8_000)
        slides.append(
            {
                "kind": "slide",
                "index": index,
                "title": title or None,
                "blocks": blocks,
                "tables": tables,
                "speaker_notes": notes or None,
            }
        )
    core = presentation.core_properties
    metadata = {
        "title": budget.take(core.title, 1_000),
        "subject": budget.take(core.subject, 1_000),
        "author": budget.take(core.author, 1_000),
        "keywords": budget.take(core.keywords, 1_000),
        "comments": budget.take(core.comments, 2_000),
        "created": _json_default(core.created) if core.created else None,
        "modified": _json_default(core.modified) if core.modified else None,
    }
    return {
        "document_type": "powerpoint",
        "metadata": {
            key: value for key, value in metadata.items() if value not in {None, ""}
        },
        "slide_count": len(slide_list),
        "slides": slides,
        "truncated": len(slide_list) > len(slides) or budget.truncated,
    }


def _bounded_tables(tables: list[Any], budget: OutputBudget) -> list[dict[str, Any]]:
    output = []
    for table in tables[:50]:
        rows = []
        for row in table[:MAX_TABLE_ROWS]:
            rows.append(
                [
                    budget.take(cell, 2_000) if cell is not None else None
                    for cell in row[:MAX_TABLE_COLUMNS]
                ]
            )
        output.append(
            {
                "rows": rows,
                "row_count": len(table),
                "column_count": max((len(row) for row in table), default=0),
                "truncated": len(table) > MAX_TABLE_ROWS
                or any(len(row) > MAX_TABLE_COLUMNS for row in table),
            }
        )
        if budget.remaining_chars <= 0:
            break
    return output


def _excel_value(value: Any, budget: OutputBudget) -> Any:
    if value is None or isinstance(value, (bool, int, float)):
        return value
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    return budget.take(value, 4_000)


def _json_default(value: Any) -> str:
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    return str(value)
